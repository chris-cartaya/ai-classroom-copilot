import os
from typing import List, Tuple

from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings

try:
    from pptx import Presentation
except Exception as e:
    Presentation = None


def _init_vector_store(persist_directory: str = "./chroma_db") -> Tuple[Chroma, OllamaEmbeddings]:
    embeddings = OllamaEmbeddings(
        model="nomic-embed-text",
        base_url="http://localhost:11434"
    )
    vector_store = Chroma(
        persist_directory=persist_directory,
        embedding_function=embeddings
    )
    return vector_store, embeddings


def pptx_to_documents(file_path: str, week_title: str) -> List[Document]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Please add python-pptx to requirements and install it.")
    prs = Presentation(file_path)
    documents: List[Document] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    texts.append(text)
        content = "\\n".join(texts).strip()
        if not content:
            continue
        documents.append(
            Document(
                page_content=content,
                metadata={
                    "source": os.path.basename(file_path),
                    "module": week_title,
                    "slide": slide_index,
                }
            )
        )
    return documents


def ingest_pptx_to_chroma(file_path: str, week_title: str, persist_directory: str = "./chroma_db") -> int:
    vector_store, _ = _init_vector_store(persist_directory=persist_directory)
    documents = pptx_to_documents(file_path=file_path, week_title=week_title)
    if not documents:
        return 0
    # Upsert by content+metadata; Chroma add_documents handles dedup by ids if provided. Keep simple here.
    vector_store.add_documents(documents)
    return len(documents)


