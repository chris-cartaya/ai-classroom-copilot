"""
FastAPI server for AI Classroom Co-Pilot RAG System
Provides API endpoints for file upload and question answering
"""

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import os
import shutil
import logging
from typing import List, Dict, Any
import tempfile
from pathlib import Path

# Import our existing RAG system
from generation import RetrievalAugmentedGeneration
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize FastAPI app
app = FastAPI(
    title="AI Classroom Co-Pilot API",
    description="RAG system for educational document Q&A",
    version="1.0.0"
)

# Add CORS middleware for frontend integration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # React dev server
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize RAG system
rag_system = RetrievalAugmentedGeneration()

# Supported file types for upload
SUPPORTED_EXTENSIONS = {'.pdf', '.txt', '.docx', '.doc', '.pptx'}
UPLOAD_DIR = Path("./uploaded_files")
CHROMA_DIR = Path("./chroma_db")

# Create directories if they don't exist
UPLOAD_DIR.mkdir(exist_ok=True)
CHROMA_DIR.mkdir(exist_ok=True)

class PPTXLoader:
    """Simple PowerPoint loader using python-pptx"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        """Load PowerPoint file and return documents"""
        from pptx import Presentation

        try:
            prs = Presentation(self.file_path)
            text_content = []

            # Extract text from all slides
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))

            full_text = "\n\n".join(text_content)

            if not full_text.strip():
                return [Document(page_content="No text content found in PowerPoint file.")]
            else:
                return [Document(page_content=full_text)]

        except Exception as e:
            logger.error(f"Error loading PPTX file {self.file_path}: {e}")
            return [Document(page_content=f"Error loading PowerPoint file: {str(e)}")]

def get_file_loader(file_path: Path):
    """Get appropriate document loader based on file extension"""
    extension = file_path.suffix.lower()

    if extension == '.pdf':
        return PyPDFLoader(str(file_path))
    elif extension in ['.txt']:
        return TextLoader(str(file_path))
    elif extension in ['.docx', '.doc']:
        return Docx2txtLoader(str(file_path))
    elif extension == '.pptx':
        return PPTXLoader(str(file_path))
    else:
        raise ValueError(f"Unsupported file type: {extension}")

def process_uploaded_file(file_path: Path, filename: str) -> List[Document]:
    """Process uploaded file and return document chunks"""
    try:
        # Load document based on file type
        loader = get_file_loader(file_path)
        documents = loader.load()

        # Split documents into chunks for better retrieval
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

        chunks = text_splitter.split_documents(documents)

        # Add metadata to chunks
        for i, chunk in enumerate(chunks):
            chunk.metadata.update({
                'source': filename,
                'chunk_index': i,
            })

        logger.info(f"Processed {filename}: {len(chunks)} chunks created")
        return chunks

    except Exception as e:
        logger.error(f"Error processing file {filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/")
async def root():
    """Root endpoint"""
    return {"message": "AI Classroom Co-Pilot API", "status": "running"}

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "upload_dir": str(UPLOAD_DIR),
        "chroma_dir": str(CHROMA_DIR)
    }

@app.post("/ask")
async def ask_question(question: str = Form(...)):
    """
    Ask a question to the RAG system
    Returns answer with citations from uploaded documents
    """
    try:
        if not question.strip():
            raise HTTPException(status_code=400, detail="Question cannot be empty")

        logger.info(f"Received question: {question}")

        # Use RAG system to get answer
        result = rag_system.ask_question(question)

        # Extract citations from retrieved documents
        citations = [{"source": doc.metadata.get('source', 'Unknown')} for doc, score in result.get('documents', [])]

        return {
            "question": result["question"],
            "answer": result["answer"],
            "documents_retrieved": result["documents_retrieved"],
            "citations": citations,
            "status": "success"
        }

    except Exception as e:
        logger.error(f"Error processing question: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing question: {str(e)}")

@app.post("/chat")
async def chat_direct(question: str = Form(...)):
    """
    Direct chat with LLM (bypasses RAG system)
    For testing basic LLM communication without document retrieval
    """
    try:
        if not question.strip():
            raise HTTPException(status_code=400, detail="Question cannot be empty")

        logger.info(f"Direct chat question: {question}")

        # Direct LLM call without retrieval
        import ollama
        response = ollama.generate(
            model="phi3:3.8b",
            prompt=question,
            options={
                'temperature': 0.7,
                'top_p': 0.9,
            }
        )

        return {
            "question": question,
            "answer": response['response'],
            "model": "phi3:3.8b",
            "mode": "direct_llm",
            "status": "success"
        }

    except Exception as e:
        logger.error(f"Error in direct chat: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing chat: {str(e)}")

@app.post("/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    """
    Upload files to the RAG system
    Files are processed and added to the vector database
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    uploaded_files = []
    failed_files = []

    try:
        for file in files:
            # Validate file extension
            file_extension = Path(file.filename).suffix.lower()
            if file_extension not in SUPPORTED_EXTENSIONS:
                failed_files.append({
                    "filename": file.filename,
                    "error": f"Unsupported file type: {file_extension}"
                })
                continue

            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
                # Write uploaded file content to temp file
                content = await file.read()
                temp_file.write(content)
                temp_file_path = Path(temp_file.name)

            try:
                # Process the file
                chunks = process_uploaded_file(temp_file_path, file.filename)

                # Add chunks to vector database
                if chunks:
                    # Import here to avoid circular imports
                    from retrieval import SlideRetriever
                    retriever = SlideRetriever()

                    # Add documents to ChromaDB
                    retriever.vector_store.add_documents(chunks)

                    uploaded_files.append({
                        "filename": file.filename,
                        "chunks_created": len(chunks),
                        "status": "success"
                    })
                else:
                    failed_files.append({
                        "filename": file.filename,
                        "error": "No content could be extracted"
                    })

            except Exception as e:
                logger.error(f"Error processing {file.filename}: {e}")
                failed_files.append({
                    "filename": file.filename,
                    "error": str(e)
                })
            finally:
                # Clean up temp file
                if temp_file_path.exists():
                    temp_file_path.unlink()

        return {
            "uploaded_files": uploaded_files,
            "failed_files": failed_files,
            "total_uploaded": len(uploaded_files),
            "total_failed": len(failed_files)
        }

    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.get("/documents")
async def list_documents():
    """List all documents currently stored in the vector database"""
    try:
        from retrieval import SlideRetriever
        retriever = SlideRetriever()

        # Get all documents from the database
        all_docs = retriever.vector_store.get()

        if not all_docs or not all_docs['ids']:
            return {
                "total_chunks": 0,
                "sources": [],
                "documents": [],
                "status": "empty"
            }

        # Group documents by source
        sources = {}
        for metadata in all_docs['metadatas']:
            source = metadata.get('source', 'Unknown')
            if source not in sources:
                sources[source] = 0
            sources[source] += 1

        # Prepare document list
        documents = []
        for i, (doc_id, metadata) in enumerate(zip(all_docs['ids'], all_docs['metadatas'])):
            documents.append({
                "id": doc_id,
                "source": metadata.get('source', 'Unknown'),
                "chunk_index": metadata.get('chunk_index', 0),
                "content_preview": all_docs['documents'][i][:200] + "..." if len(all_docs['documents'][i]) > 200 else all_docs['documents'][i]
            })

        return {
            "total_chunks": len(all_docs['ids']),
            "total_sources": len(sources),
            "sources": sources,
            "documents": documents[:50],  # Limit to first 50 for performance
            "status": "success"
        }

    except Exception as e:
        logger.error(f"Error listing documents: {e}")
        raise HTTPException(status_code=500, detail=f"Error listing documents: {str(e)}")

@app.delete("/documents")
async def clear_documents():
    """Clear all documents from the vector database"""
    try:
        from retrieval import SlideRetriever
        retriever = SlideRetriever()

        # Get all document IDs and delete them
        all_docs = retriever.vector_store.get()
        if all_docs and all_docs['ids']:
            retriever.vector_store.delete(all_docs['ids'])

        return {"message": "All documents cleared from database", "status": "success"}

    except Exception as e:
        logger.error(f"Error clearing documents: {e}")
        raise HTTPException(status_code=500, detail=f"Error clearing documents: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
